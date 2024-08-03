import os
from pypdf import PdfReader
from pptx import Presentation
import pathlib

class Parser:
    def __init__(self, input_path):
        if not os.path.exists(input_path):
            raise FileExistsError(f"{input_path} does not exist")
        
        self.input_path = os.path.abspath(input_path)
        self.file_name, self.file_ext = os.path.splitext(self.input_path)
        
    def get_output_path_from_input(self, output_file_name = None, output_file_ext = ".txt"):
        output_file_name = pathlib.Path(self.input_path if output_file_name == None else output_file_name).stem + output_file_ext
        output_path = os.path.join(os.path.dirname(self.input_path), output_file_name)
    
        return output_path
        
    def convertToText(self) -> str:
        if self.file_ext == ".pptx":
            return self._pptToText()
        elif self.file_ext == ".pdf":
            return self._pdfToText()
        elif self.file_ext == ".txt":
            return self._txtToText()
        else:
            raise NotImplementedError("Only pdf & ppt required")
        
    def _pptToText(self) -> str:
        presentation = Presentation(self.input_path)
    
        text = ""
    
        for slide in presentation.slides: 
            for shape in slide.shapes: 
                if hasattr(shape, "text"): 
                    text += shape.text
        return text
    
    def _pdfToText(self) -> str:
        pdf = PdfReader(self.input_path)
        
        text = "".join([page.extract_text() for page in pdf.pages])
        
        return text
    
    def _txtToText(self) -> str:
        text = ""
        with open(self.input_path, 'r') as f:
            lines = f.readlines()
            text = "\n".join(lines)
            
        return text
        
    def convertToTextAndSave(self):
        text = self.convertToText()
        
        output_path = self.get_output_path_from_input(output_file_ext= ".txt")
        with open(output_path, "w") as f:
            f.write(text)
        
        print(f"Converted and save to {output_path}")